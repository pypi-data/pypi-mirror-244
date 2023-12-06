"""
    PPT版式(宽:高):
        - "4:3" : "25.4cm:19.05cm, 10英寸:7.5英寸",
        - "16:9": "33.87cm:19.05cm, 10英寸:7.5英寸",
        - "A4"  : "29.7cm:21cm, 11.693英寸:8.268英寸",
"""
import pathlib
from itertools import product

import pandas as pd
from pptx import Presentation, slide
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION, XL_MARKER_STYLE, XL_LABEL_POSITION

from automate_office.automate_ppt.ppt_enums import PARAGRAPH_ALIGN, DEFAULT_PARAGRAPH_SPACING
from automate_office.automate_ppt.Charts import ChartCreator, set_color, set_font


class PPTCreator(object):

    def __init__(self, size_params=None):
        self._create_pptx(size_params)

    def _get_size(self, size_params):
        """
            size_params: 字典, {"unit_type": "inches", "width": 10, "height": 5.625}
        """
        if size_params is None:
            size_params = {
                "unit_type": "inches", 
                "width": 40/3, 
                "height": 7.5
            }
        if isinstance(size_params, dict):
            unit_type = {
                "inches": Inches,
                "cm": Cm
            }.get(
                size_params.get("unit_type", "inches"), Inches
            )

            self.height = unit_type(size_params.get("height", 7.5))
            self.width = unit_type(size_params.get("width", 40/3))

        else:
            raise ValueError("size_params要求是字典格式")

    @staticmethod
    def get_absolute_cordinate(length, rel_ratio):
        return Inches(length.inches * rel_ratio)

    def _create_pptx(self, size_params):
        self._get_size(size_params=size_params)
        self.prs = Presentation()
        self.prs.slide_height = self.height
        self.prs.slide_width = self.width

    def get_coordinate(self, rel_coordinate):
        rel_left, rel_top, rel_width, rel_height = rel_coordinate
        left = self.get_absolute_cordinate(length=self.width, rel_ratio=rel_left)
        top = self.get_absolute_cordinate(length=self.height, rel_ratio=rel_top)
        width = self.get_absolute_cordinate(length=self.width, rel_ratio=rel_width)
        height = self.get_absolute_cordinate(length=self.height, rel_ratio=rel_height)
        return left, top, width, height
    
    def create_slide(self, slide_idx=6):
        """
            0: "Title (presentation title slide)",
            1: "Title and Content",
            2: "Section Header (sometimes called Segue)",
            3: "Two Content (side by side bullet textboxes)",
            4: "Comparison (same but additional title for each side by side content box)",
            5: "Title Only",
            6: "Blank",
            7: "Content with Caption",
            8: "Picture with Caption"
        """ 
        slide_layout = self.prs.slide_layouts[slide_idx]
        return self.prs.slides.add_slide(slide_layout)

    def save_pptx(self, filename=None, save_path=None):
        if save_path is None:
            save_path = pathlib.Path.home().joinpath("Downloads")
        elif not isinstance(save_path, pathlib.WindowsPath):
            try:
                save_path = pathlib.Path(save_path)
            except Exception as e:
                raise ValueError(f"存储路径错误，{e}")

            if not save_path.is_dir():
                raise ValueError(f"{save_path}不是目录")

        if filename is None:
            filename = "output.pptx"

        self.prs.save(save_path.joinpath(filename))

    def add_textbox(
            self, 
            cur_slide: slide.Slide, 
            rel_coordinate=(0, 0, 1, 1), 
            paragraph_infos=None,
            vertical_anchor = MSO_ANCHOR.MIDDLE
        ):
        left, top, width, height = self.get_coordinate(rel_coordinate=rel_coordinate)
        shape = cur_slide.shapes.add_textbox(
            left, top, width, height
        )

        text_frame = shape.text_frame
        # 自动换行
        text_frame.word_wrap = True
        # 垂直居中, TOP MIDDLE BOTTOM
        text_frame.vertical_anchor = vertical_anchor
        if paragraph_infos is not None:
            self.add_text(text_frame=text_frame, paragraph_infos=paragraph_infos)
    
    def add_text(self, text_frame, paragraph_infos):
        first_paragraph = True
        for paragraph_info in paragraph_infos:
            if first_paragraph:
                p = text_frame.paragraphs[0]
                first_paragraph = False
            else:
                p = text_frame.add_paragraph()
            
            p.alignment = PARAGRAPH_ALIGN.get(paragraph_info.get("对齐", "左对齐"), PP_ALIGN.LEFT)
            p.space_after = Pt(paragraph_info.get("段落间距", DEFAULT_PARAGRAPH_SPACING))
            p.level = paragraph_info.get("level", 0)

            texts_infos = paragraph_info.get("文字", [])
            if texts_infos:
                for text, font_info in texts_infos:
                    run = p.add_run()
                    run.text = text
                    set_font(run.font, font_info)
    
    def add_rectangle_shape(
        self, 
        cur_slide: slide.Slide, 
        rel_coordinate=(0, 0, 1, 1), 
        paragraph_infos=None,
        vertical_anchor=MSO_ANCHOR.TOP,
        background_color=(255, 255, 255),
        line_color=(255, 255, 255),
        rounded=False
    ):  
        """可以用下面这个代码看形状的具体样子
                import pathlib
                from pptx import Presentation
                from pptx.util import Inches
                from pptx.enum.shapes import MSO_SHAPE

                prs = Presentation()
                slide_layout = prs.slide_layouts[6]
                for shape_name in [i for i in dir(MSO_SHAPE) if not i.startswith("_")]:
                    slide = prs.slides.add_slide(slide_layout)
                    left, top, width, height = Inches(1), Inches(1), Inches(3), Inches(3)
                    shape_obj = getattr(MSO_SHAPE, shape_name)
                    try:
                        shape = slide.shapes.add_shape(
                            shape_obj, left, top, width, height
                        )
                        text_frame = shape.text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = shape_name
                    except Exception as e:
                        # 删除这个slide
                        xml_slides = prs.slides._sldIdLst  # pylint: disable=W0212
                        slides = list(xml_slides)
                        xml_slides.remove(slides[-1])
                save_path = pathlib.Path.home().joinpath("Downloads")
                prs.save(save_path.joinpath("shapes.pptx"))
        """
        left, top, width, height = self.get_coordinate(rel_coordinate=rel_coordinate)
        shape_name = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = cur_slide.shapes.add_shape(
            shape_name, left, top, width, height
        )
        # 填充颜色
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*background_color)
        # 轮廓颜色
        line = shape.line
        line.color.rgb = RGBColor(*line_color)
        # 取消阴影
        shape.shadow.inherit = False

        # 填充文字
        text_frame = shape.text_frame
        # 自动换行
        text_frame.word_wrap = True
        # 垂直居中, TOP MIDDLE BOTTOM
        text_frame.vertical_anchor = vertical_anchor
        if paragraph_infos is not None:
            self.add_text(text_frame=text_frame, paragraph_infos=paragraph_infos)


    def add_picture(self, cur_slide: slide.Slide, image_path, rel_coordinate=(0, 0, 1, 1)):
        left, top, width, height = self.get_coordinate(rel_coordinate=rel_coordinate)
        _ = cur_slide.shapes.add_picture(
            image_file=image_path, left=left, top=top, width=width, height=height
        )
    
    def add_bar_chart(
        self,
        cur_slide: slide.Slide, 
        chart_data: CategoryChartData,
        rel_coordinate=(0, 0, 1, 1),
        bar_type="纵向",
        title=None,
        title_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
        series_colors=None,
        tick_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
        rotation_x=None,
        rotation_y=None,
        data_labels_font_info=None,
        data_labels_position=XL_LABEL_POSITION.INSIDE_END,
        data_labels_format=None
    ):
        ChartCreator().create_bar_chart(
            cur_slide=cur_slide, 
            chart_data=chart_data,
            absolute_coordinate=self.get_coordinate(rel_coordinate=rel_coordinate),
            bar_type=bar_type,
            title=title,
            title_font_info=title_font_info,
            series_colors=series_colors,
            tick_font_info=tick_font_info,
            rotation_x=rotation_x,
            rotation_y=rotation_y,
            data_labels_font_info=data_labels_font_info,
            data_labels_position=data_labels_position,
            data_labels_format=data_labels_format
        )

    def add_pie_chart(
            self, 
            cur_slide: slide.Slide, 
            chart_data: CategoryChartData,
            rel_coordinate=(0, 0, 1, 1),
            title=None,
            title_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
            series_colors=None,
            data_labels_font_info=None,
            data_labels_position=XL_LABEL_POSITION.OUTSIDE_END,
            legend_font_info=None,
            legend_position=XL_LEGEND_POSITION.BOTTOM, 
            data_labels_format="0.00%"
        ):
        ChartCreator().create_pie_chart(
            cur_slide=cur_slide, 
            chart_data=chart_data,
            absolute_coordinate=self.get_coordinate(rel_coordinate=rel_coordinate),
            title=title,
            title_font_info=title_font_info,
            series_colors=series_colors,
            data_labels_font_info=data_labels_font_info,
            data_labels_position=data_labels_position,
            legend_font_info=legend_font_info,
            legend_position=legend_position, 
            data_labels_format=data_labels_format
        )
    
    def add_line_chart(
        self,
        cur_slide: slide.Slide,
        chart_data: CategoryChartData,
        rel_coordinate=(0, 0, 1, 1),
        title='标题',
        title_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
        legend_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
        legend_position=XL_LEGEND_POSITION.BOTTOM,
        series_colors=None,
        marker_colors=None,
        marker_style=XL_MARKER_STYLE.CIRCLE,
        smooth=False,
        tick_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
        axis_color=(0, 0, 0),
        rotation_x=None,
        rotation_y=None,
        data_labels_font_info=None,
        data_labels_position=XL_LABEL_POSITION.INSIDE_END
    ):
        ChartCreator().create_line_chart(
            cur_slide=cur_slide, 
            chart_data=chart_data,
            absolute_coordinate=self.get_coordinate(rel_coordinate=rel_coordinate),
            title=title,
            title_font_info=title_font_info,
            legend_font_info=legend_font_info,
            legend_position=legend_position,
            series_colors=series_colors,
            marker_colors=marker_colors,
            marker_style=marker_style,
            smooth=smooth,
            tick_font_info=tick_font_info,
            axis_color=axis_color,
            rotation_x=rotation_x,
            rotation_y=rotation_y,
            data_labels_font_info=data_labels_font_info,
            data_labels_position=data_labels_position
        )

    def add_table(
        self, 
        cur_slide: slide.Slide, 
        data: pd.DataFrame, 
        rel_coordinate=(0, 0, 1, 1),
        col_widths=None,
        row_heights=None,
        header_font_info=("微软雅黑", 14, True, False, (255, 255, 255)),
        header_fore_color=(79, 128, 189),
        odd_cell_fore_color=(233, 237, 244),
        even_cell_fore_color=(255, 255, 255),
        cell_font_info=("微软雅黑", 12, False, False, (0, 0, 0))
    ):
        num_rows = data.shape[0] + 1
        num_cols = data.shape[1]
        
        # 表格大小
        left, top, width, height = self.get_coordinate(rel_coordinate=rel_coordinate)
        shape = cur_slide.shapes.add_table(rows=num_rows, cols=num_cols, left=left, top=top, width=width, height=height)
        table = shape.table

        # 列宽
        if col_widths is None:
            col_widths = [int(1 / num_cols * 100) / 100] * num_cols
            col_widths[-1] = 1 - sum(col_widths[:-1])
        for cur_col_idx, width_ratio in zip(range(num_cols), col_widths):
            table.columns[cur_col_idx].width = self.get_absolute_cordinate(length=width, rel_ratio=width_ratio)

        # 行高
        if row_heights is None:
            row_heights =[int(1 / num_rows * 100) / 100] * num_rows
            row_heights[-1] = 1 - sum(row_heights[:-1])
        for cur_row_idx, height_ratio in zip(range(num_rows), row_heights):
            table.rows[cur_row_idx].height = self.get_absolute_cordinate(length=height, rel_ratio=height_ratio)

        # 设置值和格式
        data_with_header = [data.columns.tolist()] + data.values.tolist()
        for cur_row_idx, cur_col_idx in product(range(num_rows), range(num_cols)):
            
            cur_cell = table.cell(cur_row_idx, cur_col_idx)
            cur_cell.text = f"{data_with_header[cur_row_idx][cur_col_idx]}"

            if cur_row_idx == 0:
                cur_cell_color = header_fore_color
                cur_cell_font_info = header_font_info
            else:
                cur_cell_font_info = cell_font_info
                cur_cell_color = odd_cell_fore_color if cur_row_idx % 2 == 1 else even_cell_fore_color
            set_color(cur_cell, cur_cell_color)
            
            cur_text_frame = cur_cell.text_frame
            cur_paragraph = cur_text_frame.paragraphs[0]
            set_font(
                cur_font=cur_paragraph.font,
                cur_font_info=cur_cell_font_info
            )
            # 水平居中
            cur_paragraph.alignment = PARAGRAPH_ALIGN["居中"]
            # 自动换行
            cur_text_frame.word_wrap = True
            # 垂直居中
            cur_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
