from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from automate_office.automate_ppt.ppt_enums import PARAGRAPH_ALIGN, DEFAULT_PARAGRAPH_SPACING
from automate_office.automate_ppt.helper_func import set_font


def set_textframe(text_frame, paragraph_infos=None, vertical_anchor = MSO_ANCHOR.MIDDLE):
    # text_frame.clear()

    # 自动换行
    text_frame.word_wrap = True
    # 垂直居中, TOP MIDDLE BOTTOM
    text_frame.vertical_anchor = vertical_anchor

    for idx, cur_paragraph_info in enumerate(paragraph_infos):
        cur_paragraph = text_frame.paragraphs[idx]
        cur_paragraph.clear()

        cur_paragraph.alignment = PARAGRAPH_ALIGN.get(cur_paragraph_info.get("对齐", "左对齐"), PP_ALIGN.LEFT)
        cur_paragraph.space_after = Pt(cur_paragraph_info.get("段落间距", DEFAULT_PARAGRAPH_SPACING))
        cur_paragraph.level = cur_paragraph_info.get("level", 0)

        texts_infos = cur_paragraph_info.get("文字", [])
        if texts_infos:
            for text, font_info in texts_infos:
                run = cur_paragraph.add_run()
                run.text = text
                set_font(run.font, font_info)


def replace_chart_data(chart, chart_data):
    chart.replace_data(chart_data)
