import warnings

from pptx import slide
from pptx.util import Inches
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData 
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE, XL_LABEL_POSITION

from automate_office.automate_ppt.helper_func import set_font, set_color, RGBColor


class ChartCreator(object):

    @staticmethod
    def set_title(cur_chart: Chart, title='标题', cur_font_info=("微软雅黑", 18, False, False, (89, 89, 89))):
        cur_chart.has_title = True
        cur_chart.chart_title.has_text_frame = True
        title_paragraph = cur_chart.chart_title.text_frame.paragraphs[0]
        title_paragraph.text = title
        set_font(
            cur_font=title_paragraph.font,
            cur_font_info=cur_font_info
        )
    
    @staticmethod
    def set_legend(
        cur_chart: Chart, 
        legend_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
        legend_position=XL_LEGEND_POSITION.BOTTOM
    ):
        # 设置图例
        cur_chart.has_legend = True
        cur_chart.legend.position = legend_position
        cur_chart.legend.include_in_layout = False
        set_font(
            cur_font=cur_chart.legend.font,
            cur_font_info=legend_font_info
        )

    @staticmethod
    def set_series_color(cur_chart: Chart, series_colors, color_type='self'):
        """

        Args:
            cur_chart (Chart): _description_
            series_colors (_type_): _description_
            color_type (str, optional): 在折线图里，这个就是line
        """
        try:
            plot = cur_chart.plots[0]
            for series, color in zip(plot.series, series_colors):
                if color_type != "self":
                    color_obj = getattr(series.format, color_type)
                else:
                    color_obj = series.format
                set_color(color_obj, color)
        except Exception as e:
            warnings.warn(f"系列颜色设置失败, {e}")


    @staticmethod
    def set_marker(cur_chart: Chart, marker_colors, marker_style=XL_MARKER_STYLE.CIRCLE):
        try:
            plot = cur_chart.plots[0]
            for i, series in enumerate(plot.series):
                # 设置marker样式
                series.marker.style = marker_style

                if marker_colors is not None and i < len(marker_colors):
                    color = marker_colors[i]
                    # marker颜色
                    set_color(obj=series.marker.format, color=color)
                    # marker轮廓颜色，不设置的话就是默认颜色，和设置的颜色不一样了
                    set_color(obj=series.marker.format.line, color=color)
             
        except Exception as e:
            warnings.warn(f"图例样式设置失败, {e}")

    @staticmethod
    def set_tick(cur_chart: Chart, tick_name, tick_font_info, axis_color=(0, 0, 0), rotation=None):
        if tick_name == "x":
            # 横坐标轴
            cur_axis = cur_chart.category_axis
            # tick_labels = cur_chart.category_axis.tick_labels
        elif tick_name == "y":
            # 纵坐标轴
            cur_axis = cur_chart.value_axis
            # tick_labels = cur_chart.value_axis.tick_labels
        else:
            raise ValueError(f"轴名称出错, 输入的是{tick_name}")

        set_font(
            cur_font=cur_axis.tick_labels.font,
            cur_font_info=tick_font_info
        )
        cur_chart.category_axis.format.line.color.rgb = RGBColor(*axis_color)

        if rotation is not None:
            txPr = cur_axis.tick_labels._element.get_or_add_txPr()
            # X轴标签旋转45度
            # 在 OpenXML SDK 里面，采用的基础单位是 60000 倍的 Degree 角度值
            # 也就是在获取到 OpenXML 的 Int32Value 时，获取数值，除以 60000 就拿到了角度值
            # 例如： 逆时针旋转45度，那么就等于 60000 * (-45) = -2700000
            txPr.bodyPr.set("rot", f"{60000 * rotation:>.0f}")

    @staticmethod
    def set_data_labels(cur_chart: Chart, data_labels_font_info, data_labels_position=XL_LABEL_POSITION.INSIDE_END, data_labels_format=None):
        plot = cur_chart.plots[0]
        plot.has_data_labels = True

        for series in cur_chart.plots[0].series:
            data_labels = series.data_labels
            set_font(data_labels.font, data_labels_font_info)
            data_labels.position = data_labels_position
            data_labels.show_value = True
            if data_labels_format is not None:
                data_labels.number_format = data_labels_format

    def create_line_chart(
            self,
            cur_slide: slide.Slide, 
            chart_data: CategoryChartData,
            absolute_coordinate=(Inches(0), Inches(0), Inches(7.5), Inches(7.5)),
            title=None,
            title_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
            legend_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
            legend_position=XL_LEGEND_POSITION.BOTTOM,
            series_colors=None,
            marker_colors=None,
            marker_style=XL_MARKER_STYLE.CIRCLE,
            smooth=False,
            tick_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
            axis_color=(0, 0, 0),
            rotation_x=None,
            rotation_y=None,
            data_labels_font_info=None,
            data_labels_position=XL_LABEL_POSITION.INSIDE_END
        ):
        x, y, cx, cy = absolute_coordinate
        graphic_frame = cur_slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart

        # 设置标题
        if title is not None:
            self.set_title(cur_chart=chart, title=title, cur_font_info=title_font_info)
        # 设置图例
        self.set_legend(cur_chart=chart, legend_font_info=legend_font_info, legend_position=legend_position)
        # 设置系列颜色
        if series_colors is not None:
            self.set_series_color(
                cur_chart=chart, 
                series_colors=series_colors, 
                color_type='line'
            )
        # 设置点样式
        self.set_marker(
            cur_chart=chart, 
            marker_colors=marker_colors, 
            marker_style=marker_style
        )

        # 去掉网格
        chart.value_axis.has_major_gridlines = False

        # 平滑折线图
        for series in chart.series:
            series.smooth = smooth

        # 设定x轴y轴
        self.set_tick(cur_chart=chart, tick_name="x", tick_font_info=tick_font_info, axis_color=axis_color, rotation=rotation_x)
        self.set_tick(cur_chart=chart, tick_name="y", tick_font_info=tick_font_info, axis_color=axis_color, rotation=rotation_y)

        # 设定数据标注data_labels
        if data_labels_font_info is not None:
            self.set_data_labels(cur_chart=chart, data_labels_font_info=data_labels_font_info, data_labels_position=data_labels_position)

    def create_bar_chart(
            self,
            cur_slide: slide.Slide, 
            chart_data: CategoryChartData,
            absolute_coordinate=(Inches(0), Inches(0), Inches(7.5), Inches(7.5)),
            bar_type="纵向",
            title=None,
            title_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
            series_colors=None,
            tick_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
            axis_color=(0, 0, 0),
            rotation_x=None,
            rotation_y=None,
            data_labels_font_info=None,
            data_labels_position=XL_LABEL_POSITION.INSIDE_END, 
            data_labels_format=None
        ):
        x, y, cx, cy = absolute_coordinate
        if bar_type == "纵向":
            chart_name = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif bar_type == "横向":
            chart_name = XL_CHART_TYPE.BAR_CLUSTERED
        else:
            chart_name = XL_CHART_TYPE.COLUMN_CLUSTERED
        graphic_frame = cur_slide.shapes.add_chart(
            chart_name, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart

        # 设置标题
        if title is not None:
            self.set_title(cur_chart=chart, title=title, cur_font_info=title_font_info)

        # 去掉网格线
        chart.value_axis.has_major_gridlines = False

        # 设置系列颜色
        if series_colors is not None:
            self.set_series_color(
                cur_chart=chart, 
                series_colors=series_colors, 
                color_type="self"
            )

        # 设定x轴y轴
        self.set_tick(cur_chart=chart, tick_name="x", tick_font_info=tick_font_info, axis_color=axis_color, rotation=rotation_x)
        self.set_tick(cur_chart=chart, tick_name="y", tick_font_info=tick_font_info, axis_color=axis_color, rotation=rotation_y)

        # 设定数据标注data_labels
        if data_labels_font_info is not None:
            self.set_data_labels(cur_chart=chart, data_labels_font_info=data_labels_font_info, data_labels_position=data_labels_position, data_labels_format=data_labels_format)

    def create_pie_chart(
            self,
            cur_slide: slide.Slide, 
            chart_data: CategoryChartData,
            absolute_coordinate=(Inches(0), Inches(0), Inches(7.5), Inches(7.5)),
            title=None,
            title_font_info=("微软雅黑", 18, False, False, (89, 89, 89)),
            series_colors=None,
            data_labels_font_info=None,
            data_labels_position=XL_LABEL_POSITION.OUTSIDE_END,
            legend_font_info=None,
            legend_position=XL_LEGEND_POSITION.BOTTOM, 
            data_labels_format="0.00%"
        ):

        x, y, cx, cy = absolute_coordinate

        graphic_frame = cur_slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        )
        chart = graphic_frame.chart

        # 设置标题
        if title is not None:
            self.set_title(cur_chart=chart, title=title, cur_font_info=title_font_info)

        # # 去掉网格线
        # chart.value_axis.has_major_gridlines = False

        # 设置系列颜色
        if series_colors is not None:
            for point, color in zip(chart.series[0].points, series_colors):
                set_color(point.format, color)

        # 设定数据标注data_labels
        if data_labels_font_info is not None:
            self.set_data_labels(cur_chart=chart, data_labels_font_info=data_labels_font_info, data_labels_position=data_labels_position, data_labels_format=data_labels_format)

        # 设置图例
        if legend_font_info is not None:
            chart.has_legend = True
            chart.legend.include_in_layout = False
            self.set_legend(cur_chart=chart, legend_font_info=legend_font_info, legend_position=legend_position)



       

